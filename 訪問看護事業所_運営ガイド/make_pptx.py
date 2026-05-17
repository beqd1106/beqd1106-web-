"""障碍者福祉事業所 運営ガイド — 概要プレゼンテーション生成スクリプト"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.dml.color import RGBColor

# ── カラー定義 ──────────────────────────────────
C_CORAL   = RGBColor(0xD4, 0x87, 0x5A)   # テラコッタ（メイン）
C_DEEP    = RGBColor(0x7A, 0x30, 0x18)   # 深いブラウンレッド
C_SAGE    = RGBColor(0x7A, 0xB5, 0x92)   # セージグリーン
C_AMBER   = RGBColor(0xF0, 0xB8, 0x40)   # アンバー
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_CREAM   = RGBColor(0xFF, 0xFB, 0xF7)
C_DARK    = RGBColor(0x3A, 0x2D, 0x24)   # 暖色ダーク
C_MID     = RGBColor(0x7A, 0x65, 0x58)   # ミッドブラウン
C_LIGHT_B = RGBColor(0xFA, 0xEA, 0xDA)   # 薄ピーチ

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "障碍者福祉事業所_運営ガイド_概要資料.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 完全ブランク

# ── ヘルパー関数 ───────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w: shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True):
    txbox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Meiryo UI"
    if color: run.font.color.rgb = color
    return txbox

def add_para(txbox, text, size=11, bold=False, color=None, align=PP_ALIGN.LEFT, space_before=0):
    tf = txbox.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Meiryo UI"
    if color: run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None):
    """スライド共通ヘッダー帯"""
    add_rect(slide, 0, 0, 13.33, 1.4, fill=C_DEEP)
    add_rect(slide, 0, 1.35, 13.33, 0.06, fill=C_CORAL)
    add_text(slide, title, 0.4, 0.1, 12, 0.8, size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.45, 0.85, 11, 0.45, size=13, color=C_CORAL, bold=False)
    add_rect(slide, 0, 7.35, 13.33, 0.15, fill=C_CORAL)
    add_text(slide, "障碍者福祉事業所 運営ガイド  |  京都府対応版  2026",
             0.3, 7.3, 10, 0.2, size=8, color=C_WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — タイトル
# ═══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.33, 7.5, fill=C_CREAM)
add_rect(s1, 0, 0, 13.33, 4.0, fill=C_DEEP)
add_rect(s1, 0, 3.95, 13.33, 0.1, fill=C_CORAL)

# タイトル
add_text(s1, "障碍者福祉事業所", 0.5, 0.5, 12, 0.9, size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "運 営 ガ イ ド", 0.5, 1.35, 12, 0.9, size=44, bold=True, color=C_CORAL, align=PP_ALIGN.CENTER)
add_text(s1, "訪問看護  ／  就労継続支援A型・B型  ／  相談支援事業所",
         0.5, 2.3, 12, 0.5, size=16, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "京都府・京都市 対応版  ｜  2026年度",
         0.5, 2.85, 12, 0.4, size=13, color=C_LIGHT_B, align=PP_ALIGN.CENTER)

# 概要ボックス（下半分）
boxes = [
    ("13", "ページ\nのHTMLガイド", C_CORAL),
    ("56", "種類の\n様式フォーマット", C_SAGE),
    ("Q&A", "＋年間\nスケジュール", C_AMBER),
    ("無料", "で開ける\nブラウザ完結", RGBColor(0x9E,0x4A,0x28)),
]
for i, (num, label, col) in enumerate(boxes):
    x = 0.6 + i * 3.1
    add_rect(s1, x, 4.3, 2.8, 2.6, fill=C_WHITE)
    add_rect(s1, x, 4.3, 2.8, 0.18, fill=col)
    add_text(s1, num,   x+0.1, 4.45, 2.6, 1.1, size=48, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s1, label, x+0.1, 5.55, 2.6, 0.9, size=13, color=C_MID, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — このサイトが解決すること
# ═══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.33, 7.5, fill=C_CREAM)
slide_header(s2, "このガイドが解決すること", "障碍者福祉事業所の運営で直面する3大課題を一発解消")

problems = [
    ("😰", "情報が散らばっている",
     "法令・書類・申請手続き・請求方法が\nそれぞれ別の場所に存在し、\n担当者が毎回調べ直している。",
     "→ 13ページのHTMLに\n　全情報を一元集約"),
    ("📋", "書類・様式を探すのが大変",
     "指定申請書・個別支援計画・処遇改善加算\n様式など、毎回公式サイトを\n探し回っている。",
     "→ 56種の様式を\n　ボタン1クリックで取得"),
    ("📅", "期限管理が属人化している",
     "毎月10日の請求・4月15日の処遇改善\n届出・6年毎の指定更新など\n重要期限を個人が管理している。",
     "→ 年間スケジュールページで\n　全期限を一覧管理"),
]

for i, (icon, title, prob, sol) in enumerate(problems):
    x = 0.35 + i * 4.2
    add_rect(s2, x, 1.6, 3.9, 5.5, fill=C_WHITE)
    add_rect(s2, x, 1.6, 3.9, 0.15, fill=C_CORAL)

    add_text(s2, icon,  x+0.15, 1.75, 0.7, 0.6, size=28)
    add_text(s2, title, x+0.15, 1.75, 3.6, 0.6, size=15, bold=True, color=C_DEEP)

    add_text(s2, "【課題】", x+0.2, 2.45, 3.5, 0.35, size=10, bold=True, color=C_MID)
    add_text(s2, prob,  x+0.2, 2.78, 3.5, 1.5, size=11, color=C_DARK)

    add_rect(s2, x+0.15, 4.35, 3.6, 0.06, fill=C_CORAL)
    add_text(s2, "【解決】", x+0.2, 4.45, 3.5, 0.35, size=10, bold=True, color=C_SAGE)
    add_text(s2, sol,   x+0.2, 4.78, 3.5, 1.1, size=12, bold=True, color=C_DEEP)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — コンテンツマップ（13ページ）
# ═══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.33, 7.5, fill=C_CREAM)
slide_header(s3, "コンテンツマップ（全13ページ）", "index.htmlを開くだけでナビゲートできます")

pages = [
    # (タイトル, 内容, 色)
    ("⚖️ 法令・制度",       "医療系・障碍者福祉系の全法令\n京都府条例・最近の改正", C_CORAL),
    ("📋 必要書類",          "指定申請書類14種\n日常業務記録書・保管期間一覧", C_CORAL),
    ("📅 管理業務",          "月次カレンダー・スタッフ管理\n実地指導対策チェック100項目", C_CORAL),
    ("🏛️ 京都府・手続き",    "実際の電話番号・窓口URL\n申請フロー・様式ダウンロードリンク", C_CORAL),
    ("💴 訪問看護 請求業務", "療養費の仕組み・加算一覧\n月次スケジュール・返戻対処法", C_CORAL),
    ("🏢 就労継続支援A型",   "スコア方式200点の解説\n最低賃金対応・申請書類一覧", C_SAGE),
    ("🏭 就労継続支援B型",   "工賃管理・京都市総量規制\n6:1配置新設・加算一覧", C_SAGE),
    ("💬 相談支援事業所",    "計画相談・地域相談の違い\n相談支援専門員の要件・報酬", C_SAGE),
    ("🔍 監査・運営指導",    "運営指導vs監査の違い\n当日フロー・自主点検100項目", RGBColor(0x9E,0x4A,0x28)),
    ("📈 収入向上・加算",    "処遇改善加算Ⅰ〜Ⅴ最大化\n稼働率・多機能型の戦略", RGBColor(0x9E,0x4A,0x28)),
    ("📆 年間スケジュール",  "月別期限カレンダー\n研修・指定更新・複数年サイクル", RGBColor(0x9E,0x4A,0x28)),
    ("❓ よくある質問",      "5カテゴリ・20問のQ&A\n現場でよくある疑問に完全回答", RGBColor(0x9E,0x4A,0x28)),
    ("📁 様式フォーマット集","56ファイルをカテゴリ別\nワンクリックでExcel・Wordを取得", C_AMBER),
]

cols = 4
for i, (title, body, col) in enumerate(pages):
    r, c = divmod(i, cols)
    x = 0.25 + c * 3.24
    y = 1.6  + r * 1.88
    add_rect(s3, x, y, 3.1, 1.72, fill=C_WHITE)
    add_rect(s3, x, y, 3.1, 0.14, fill=col)
    add_text(s3, title, x+0.1, y+0.18, 2.9, 0.5, size=11, bold=True, color=C_DEEP)
    add_text(s3, body,  x+0.1, y+0.68, 2.9, 0.9, size=9,  color=C_MID)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — 様式フォーマット集
# ═══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.33, 7.5, fill=C_CREAM)
slide_header(s4, "様式フォーマット集（56ファイル）", "公式様式23件＋オリジナル33件をブラウザから直接ダウンロード")

cats = [
    ("01\n指定申請様式", "11", "京都府・京都市の公式申請書\n事前相談票・従業者一覧・資金計画\n就労継続支援・相談支援申請書　など", C_CORAL),
    ("02\n訪問看護記録", "9",  "訪問看護計画書（厚労省公式）\n指示書・記録書Ⅰ・Ⅱ・報告書\n精神科版・退院前カンファレンス記録", C_CORAL),
    ("03\n就労支援記録", "11", "個別支援計画（京都府公式）\nモニタリング・アセスメント\n工賃台帳・賃金台帳・日報　など", C_SAGE),
    ("04\n相談支援様式", "3",  "個別支援計画参考様式（厚労省）\nこども家庭庁 別表様式\nサービス担当者会議記録", C_SAGE),
    ("05\n処遇改善加算", "6",  "計画書（令和7年度・4/15提出）\n実績報告書（7/31提出）\n記入例・変更届・特別事情届", C_AMBER),
    ("06\n事故苦情管理", "10", "事故報告書・ヒヤリハット\n苦情受付・研修記録\nシフト表・業務日誌・体制届チェック", RGBColor(0x9E,0x4A,0x28)),
    ("07\n運営規程ひな形","6", "訪問看護・就労継続支援B型\n特定相談支援事業所 運営規程\n重要事項説明書・秘密保持誓約書", RGBColor(0x5A,0x20,0x10)),
]

for i, (cat, cnt, desc, col) in enumerate(cats):
    r, c2 = divmod(i, 4)
    x = 0.25 + c2 * 3.26
    y = 1.65  + r * 2.75
    add_rect(s4, x, y, 3.08, 2.5, fill=C_WHITE)
    add_rect(s4, x, y, 3.08, 0.55, fill=col)
    add_text(s4, cat, x+0.08, y+0.02, 1.6, 0.5, size=10, bold=True, color=C_WHITE)
    add_text(s4, f"{cnt}件", x+2.1, y+0.05, 0.9, 0.45, size=22, bold=True, color=C_WHITE, align=PP_ALIGN.RIGHT)
    add_text(s4, desc, x+0.1, y+0.6, 2.9, 1.75, size=9, color=C_MID)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — 運営への貢献（4つの価値）
# ═══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.33, 7.5, fill=C_CREAM)
slide_header(s5, "運営への4つの貢献", "このガイドが事業所にもたらす具体的な価値")

values = [
    ("⏱️", "業務時間の短縮",
     "毎回Webで調べていた法令・申請書類・\n加算要件の検索がゼロに。\n\n"
     "▶ 様式DLは1クリック\n▶ 年間期限は一覧で確認\n▶ Q&Aで疑問を即解消",
     C_CORAL, "想定削減：月10〜20時間の検索時間"),
    ("⚖️", "コンプライアンス強化",
     "法令・運営基準・実地指導ポイントを\n正確に把握して運営できる。\n\n"
     "▶ 自主点検100項目チェック\n▶ 年間提出期限を見える化\n▶ 不正リスクの事前回避",
     C_SAGE, "実地指導での「不備ゼロ」を目指す"),
    ("💴", "収入の最大化",
     "算定できる加算の取り漏らしを防止。\n処遇改善加算Ⅰ取得で\n職員採用力も同時に向上。\n\n"
     "▶ 全加算の要件を一覧で確認\n▶ 体制届チェックシートで管理",
     C_AMBER, "処遇改善加算Ⅰ取得で+15〜20%収入増"),
    ("👥", "チーム力の向上",
     "情報を個人が抱え込まず\nチーム全員でアクセスできる\n運営ガイドとして活用。\n\n"
     "▶ ひな形で書類作成を標準化\n▶ 新人教育資料としても活用可能",
     C_DEEP, "属人化リスクを組織的に解消"),
]

for i, (icon, title, body, col, note) in enumerate(values):
    r2, c3 = divmod(i, 2)
    x = 0.3  + c3 * 6.5
    y = 1.65 + r2 * 2.85
    add_rect(s5, x, y, 6.2, 2.65, fill=C_WHITE)
    add_rect(s5, x, y, 6.2, 0.14, fill=col)
    add_text(s5, f"{icon} {title}", x+0.15, y+0.2, 5.9, 0.55, size=16, bold=True, color=C_DEEP)
    add_text(s5, body, x+0.2, y+0.75, 3.5, 1.8, size=10, color=C_DARK)
    # ノートバッジ
    add_rect(s5, x+3.75, y+0.75, 2.3, 1.0, fill=C_LIGHT_B)
    add_text(s5, note,  x+3.85, y+0.85, 2.1, 0.85, size=9, bold=True, color=col)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — 使い方ガイド
# ═══════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.33, 7.5, fill=C_CREAM)
slide_header(s6, "使い方ガイド", "ダウンロードフォルダを開くだけで使えます")

steps = [
    ("Step 1", "📁 フォルダを開く",
     "Downloads →\n訪問看護事業所_運営ガイド"),
    ("Step 2", "🌐 index.html を開く",
     "ブラウザ（Chrome等）で\nダブルクリック"),
    ("Step 3", "📖 目的のページへ",
     "ナビまたはカードから\n13ページに直接アクセス"),
    ("Step 4", "📥 様式をDL",
     "ページ内の 📥様式DL ボタン\nまたは📁様式DLナビから"),
]

for i, (step, title, body) in enumerate(steps):
    x = 0.3 + i * 3.25
    # 矢印
    if i < 3:
        add_text(s6, "→", x+2.9, 3.3, 0.5, 0.5, size=24, bold=True, color=C_CORAL, align=PP_ALIGN.CENTER)
    # カード
    add_rect(s6, x, 1.7, 3.0, 3.4, fill=C_WHITE)
    add_rect(s6, x, 1.7, 3.0, 0.55, fill=C_CORAL)
    add_text(s6, step, x+0.1, 1.72, 2.8, 0.5, size=12, bold=True, color=C_WHITE)
    add_text(s6, title, x+0.1, 2.3, 2.8, 0.65, size=16, bold=True, color=C_DEEP)
    add_text(s6, body, x+0.1, 2.95, 2.8, 1.5, size=13, color=C_MID)

# tips
add_rect(s6, 0.3, 5.3, 12.7, 1.85, fill=C_LIGHT_B)
add_rect(s6, 0.3, 5.3, 12.7, 0.1, fill=C_CORAL)
add_text(s6, "💡 活用ヒント", 0.5, 5.4, 3, 0.35, size=12, bold=True, color=C_CORAL)
tips_text = ("• スマートフォンからも閲覧OK（レスポンシブデザイン対応）\n"
             "• 印刷したい場合は各ページのChrome印刷機能を使用（Ctrl+P）\n"
             "• フォーマット集のExcelファイルはそのまま編集・保存して使用可能\n"
             "• 公式様式は改定時に再ダウンロードしてください（厚労省・京都府サイトへのリンクあり）")
add_text(s6, tips_text, 0.5, 5.72, 12.3, 1.3, size=10, color=C_DARK)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — まとめ
# ═══════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, 13.33, 7.5, fill=C_DEEP)
add_rect(s7, 0, 0, 13.33, 0.15, fill=C_CORAL)
add_rect(s7, 0, 7.35, 13.33, 0.15, fill=C_CORAL)

add_text(s7, "まとめ", 0.5, 0.4, 12, 0.7, size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s7, "このガイドがあれば、障碍者福祉事業所の運営に必要な",
         0.5, 1.1, 12, 0.5, size=16, color=C_LIGHT_B, align=PP_ALIGN.CENTER)
add_text(s7, "「調べる」「書類を探す」「期限を確認する」が\nすべて1か所で完結します。",
         0.5, 1.55, 12, 0.75, size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

summaries = [
    ("法令・制度", "10以上の法令を\n分かりやすく整理"),
    ("申請・書類", "公式様式23件を\n直接DL可能"),
    ("管理業務", "月次・年次業務を\nカレンダーで把握"),
    ("収入向上", "全加算・処遇改善の\n要件を完全網羅"),
    ("監査対策", "自主点検100項目\nチェックリスト"),
    ("Q&A", "現場の疑問に\n即答する20問"),
]

for i, (title, body) in enumerate(summaries):
    r3, c4 = divmod(i, 3)
    x = 0.5 + c4 * 4.1
    y = 2.7 + r3 * 2.1
    add_rect(s7, x, y, 3.8, 1.85, fill=RGBColor(0x5A,0x30,0x18))
    add_rect(s7, x, y, 3.8, 0.12, fill=C_CORAL)
    add_text(s7, title, x+0.15, y+0.18, 3.5, 0.45, size=14, bold=True, color=C_CORAL)
    add_text(s7, body,  x+0.15, y+0.65, 3.5, 0.9, size=12, color=C_WHITE)

prs.save(OUT)
print(f"OK: {os.path.basename(OUT)}")
